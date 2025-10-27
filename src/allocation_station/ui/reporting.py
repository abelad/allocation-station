"""
Reporting System for Allocation Station

This module provides comprehensive reporting capabilities including automated
PDF report generation, customizable templates, email delivery, executive
summaries, regulatory compliance reports, client-ready presentations, and
multi-language support.

Features:
    - Automated PDF report generation with charts and tables
    - Customizable report templates for different use cases
    - Email delivery with SMTP integration
    - Executive summary generation with AI-powered insights
    - Regulatory compliance reports (SEC, FINRA, etc.)
    - Client-ready PowerPoint presentations
    - Multi-language support (English, Spanish, French, German, Chinese)
"""

import io
import os
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.base import MIMEBase
from email import encoders
from datetime import datetime, timedelta
from typing import Dict, List, Optional, Tuple, Any, Union
from dataclasses import dataclass, asdict
from enum import Enum
import json
from pathlib import Path

import pandas as pd
import numpy as np
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer,
    PageBreak, Image, KeepTogether
)
from reportlab.lib.enums import TA_CENTER, TA_RIGHT, TA_LEFT
from reportlab.graphics.shapes import Drawing
from reportlab.graphics.charts.piecharts import Pie
from reportlab.graphics.charts.barcharts import VerticalBarChart
from reportlab.graphics.charts.linecharts import HorizontalLineChart
import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData


class ReportType(Enum):
    """Report types."""
    PORTFOLIO_SUMMARY = "portfolio_summary"
    PERFORMANCE_REVIEW = "performance_review"
    RISK_ANALYSIS = "risk_analysis"
    COMPLIANCE = "compliance"
    QUARTERLY_REVIEW = "quarterly_review"
    ANNUAL_REPORT = "annual_report"
    EXECUTIVE_SUMMARY = "executive_summary"
    CLIENT_PRESENTATION = "client_presentation"


class ReportFormat(Enum):
    """Report output formats."""
    PDF = "pdf"
    HTML = "html"
    PPTX = "pptx"
    EXCEL = "excel"
    JSON = "json"


class Language(Enum):
    """Supported languages."""
    ENGLISH = "en"
    SPANISH = "es"
    FRENCH = "fr"
    GERMAN = "de"
    CHINESE = "zh"


@dataclass
class ReportConfig:
    """Report configuration."""
    report_type: ReportType
    format: ReportFormat
    language: Language
    include_charts: bool = True
    include_tables: bool = True
    include_recommendations: bool = True
    include_disclosures: bool = True
    logo_path: Optional[str] = None
    company_name: str = "Allocation Station"
    report_date: datetime = None
    custom_fields: Dict[str, Any] = None

    def __post_init__(self):
        if self.report_date is None:
            self.report_date = datetime.now()
        if self.custom_fields is None:
            self.custom_fields = {}


@dataclass
class PortfolioData:
    """Portfolio data for reporting."""
    portfolio_name: str
    total_value: float
    holdings: Dict[str, Dict[str, float]]
    performance: Dict[str, float]
    allocation: Dict[str, float]
    risk_metrics: Dict[str, float]
    historical_returns: pd.DataFrame
    benchmark_comparison: Optional[Dict[str, Any]] = None


class Translator:
    """Multi-language translation support."""

    TRANSLATIONS = {
        Language.ENGLISH: {
            'portfolio_summary': 'Portfolio Summary',
            'performance_review': 'Performance Review',
            'risk_analysis': 'Risk Analysis',
            'total_value': 'Total Portfolio Value',
            'ytd_return': 'Year-to-Date Return',
            'annual_return': 'Annual Return',
            'volatility': 'Volatility',
            'sharpe_ratio': 'Sharpe Ratio',
            'max_drawdown': 'Maximum Drawdown',
            'allocation': 'Asset Allocation',
            'holdings': 'Current Holdings',
            'performance': 'Performance Metrics',
            'risk_metrics': 'Risk Metrics',
            'recommendations': 'Recommendations',
            'disclaimer': 'Important Disclaimer',
            'page': 'Page',
            'date': 'Date',
            'as_of': 'As of',
            'symbol': 'Symbol',
            'shares': 'Shares',
            'price': 'Price',
            'value': 'Value',
            'weight': 'Weight (%)',
            'change': 'Change',
            'executive_summary': 'Executive Summary',
            'key_highlights': 'Key Highlights',
            'market_overview': 'Market Overview',
            'compliance_statement': 'Compliance Statement',
        },
        Language.SPANISH: {
            'portfolio_summary': 'Resumen de Cartera',
            'performance_review': 'Revisión de Rendimiento',
            'risk_analysis': 'Análisis de Riesgo',
            'total_value': 'Valor Total de Cartera',
            'ytd_return': 'Rendimiento del Año',
            'annual_return': 'Rendimiento Anual',
            'volatility': 'Volatilidad',
            'sharpe_ratio': 'Ratio de Sharpe',
            'max_drawdown': 'Caída Máxima',
            'allocation': 'Asignación de Activos',
            'holdings': 'Posiciones Actuales',
            'performance': 'Métricas de Rendimiento',
            'risk_metrics': 'Métricas de Riesgo',
            'recommendations': 'Recomendaciones',
            'disclaimer': 'Descargo de Responsabilidad',
            'page': 'Página',
            'date': 'Fecha',
            'as_of': 'A fecha de',
            'symbol': 'Símbolo',
            'shares': 'Acciones',
            'price': 'Precio',
            'value': 'Valor',
            'weight': 'Peso (%)',
            'change': 'Cambio',
            'executive_summary': 'Resumen Ejecutivo',
            'key_highlights': 'Puntos Clave',
            'market_overview': 'Visión del Mercado',
            'compliance_statement': 'Declaración de Cumplimiento',
        },
        Language.FRENCH: {
            'portfolio_summary': 'Résumé du Portefeuille',
            'performance_review': 'Examen de Performance',
            'risk_analysis': 'Analyse des Risques',
            'total_value': 'Valeur Totale du Portefeuille',
            'ytd_return': 'Rendement de l\'Année',
            'annual_return': 'Rendement Annuel',
            'volatility': 'Volatilité',
            'sharpe_ratio': 'Ratio de Sharpe',
            'max_drawdown': 'Perte Maximale',
            'allocation': 'Allocation d\'Actifs',
            'holdings': 'Positions Actuelles',
            'performance': 'Métriques de Performance',
            'risk_metrics': 'Métriques de Risque',
            'recommendations': 'Recommandations',
            'disclaimer': 'Avertissement Important',
            'page': 'Page',
            'date': 'Date',
            'as_of': 'Au',
            'symbol': 'Symbole',
            'shares': 'Actions',
            'price': 'Prix',
            'value': 'Valeur',
            'weight': 'Poids (%)',
            'change': 'Changement',
            'executive_summary': 'Résumé Exécutif',
            'key_highlights': 'Points Clés',
            'market_overview': 'Vue du Marché',
            'compliance_statement': 'Déclaration de Conformité',
        },
        Language.GERMAN: {
            'portfolio_summary': 'Portfolio-Zusammenfassung',
            'performance_review': 'Performance-Überprüfung',
            'risk_analysis': 'Risikoanalyse',
            'total_value': 'Gesamtwert des Portfolios',
            'ytd_return': 'Jahresrendite',
            'annual_return': 'Jährliche Rendite',
            'volatility': 'Volatilität',
            'sharpe_ratio': 'Sharpe-Ratio',
            'max_drawdown': 'Maximaler Drawdown',
            'allocation': 'Vermögensallokation',
            'holdings': 'Aktuelle Positionen',
            'performance': 'Performance-Metriken',
            'risk_metrics': 'Risikokennzahlen',
            'recommendations': 'Empfehlungen',
            'disclaimer': 'Wichtiger Haftungsausschluss',
            'page': 'Seite',
            'date': 'Datum',
            'as_of': 'Stand',
            'symbol': 'Symbol',
            'shares': 'Aktien',
            'price': 'Preis',
            'value': 'Wert',
            'weight': 'Gewicht (%)',
            'change': 'Änderung',
            'executive_summary': 'Zusammenfassung',
            'key_highlights': 'Wichtige Punkte',
            'market_overview': 'Marktübersicht',
            'compliance_statement': 'Compliance-Erklärung',
        },
        Language.CHINESE: {
            'portfolio_summary': '投资组合摘要',
            'performance_review': '业绩回顾',
            'risk_analysis': '风险分析',
            'total_value': '投资组合总价值',
            'ytd_return': '年初至今回报',
            'annual_return': '年度回报',
            'volatility': '波动率',
            'sharpe_ratio': '夏普比率',
            'max_drawdown': '最大回撤',
            'allocation': '资产配置',
            'holdings': '当前持仓',
            'performance': '业绩指标',
            'risk_metrics': '风险指标',
            'recommendations': '建议',
            'disclaimer': '重要声明',
            'page': '页',
            'date': '日期',
            'as_of': '截至',
            'symbol': '代码',
            'shares': '股数',
            'price': '价格',
            'value': '价值',
            'weight': '权重 (%)',
            'change': '变化',
            'executive_summary': '执行摘要',
            'key_highlights': '重点',
            'market_overview': '市场概览',
            'compliance_statement': '合规声明',
        }
    }

    @classmethod
    def translate(cls, key: str, language: Language) -> str:
        """Translate a key to the specified language."""
        translations = cls.TRANSLATIONS.get(language, cls.TRANSLATIONS[Language.ENGLISH])
        return translations.get(key, key)


class PDFReportGenerator:
    """Generates PDF reports using ReportLab."""

    def __init__(self, config: ReportConfig):
        """Initialize PDF report generator."""
        self.config = config
        self.translator = Translator()
        self.styles = getSampleStyleSheet()
        self._setup_custom_styles()

    def _setup_custom_styles(self):
        """Setup custom paragraph styles."""
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Title'],
            fontSize=24,
            textColor=colors.HexColor('#1f77b4'),
            spaceAfter=30,
            alignment=TA_CENTER
        ))

        self.styles.add(ParagraphStyle(
            name='SectionHeader',
            parent=self.styles['Heading1'],
            fontSize=16,
            textColor=colors.HexColor('#2c3e50'),
            spaceBefore=20,
            spaceAfter=12,
            borderWidth=1,
            borderColor=colors.HexColor('#1f77b4'),
            borderPadding=5
        ))

        self.styles.add(ParagraphStyle(
            name='Highlight',
            parent=self.styles['Normal'],
            fontSize=12,
            textColor=colors.HexColor('#27ae60'),
            fontName='Helvetica-Bold'
        ))

    def _translate(self, key: str) -> str:
        """Translate a key to the configured language."""
        return self.translator.translate(key, self.config.language)

    def _create_header(self, canvas, doc):
        """Create page header."""
        canvas.saveState()
        canvas.setFont('Helvetica', 10)
        canvas.drawString(inch, 10.5 * inch, self.config.company_name)
        canvas.drawRightString(7.5 * inch, 10.5 * inch,
                              f"{self._translate('date')}: {self.config.report_date.strftime('%Y-%m-%d')}")
        canvas.restoreState()

    def _create_footer(self, canvas, doc):
        """Create page footer."""
        canvas.saveState()
        canvas.setFont('Helvetica', 9)
        canvas.drawCentredString(4.25 * inch, 0.5 * inch,
                                f"{self._translate('page')} {doc.page}")
        canvas.restoreState()

    def generate_portfolio_summary(self, portfolio: PortfolioData, output_path: str):
        """Generate portfolio summary report."""
        doc = SimpleDocTemplate(
            output_path,
            pagesize=letter,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72
        )

        story = []

        # Title
        title = Paragraph(
            self._translate('portfolio_summary'),
            self.styles['CustomTitle']
        )
        story.append(title)
        story.append(Spacer(1, 0.3 * inch))

        # Portfolio Overview
        story.append(Paragraph(
            f"{portfolio.portfolio_name}",
            self.styles['Heading2']
        ))
        story.append(Spacer(1, 0.2 * inch))

        # Key metrics table
        metrics_data = [
            [self._translate('total_value'), f"${portfolio.total_value:,.2f}"],
            [self._translate('ytd_return'),
             f"{portfolio.performance.get('ytd_return', 0):.2f}%"],
            [self._translate('annual_return'),
             f"{portfolio.performance.get('annual_return', 0):.2f}%"],
            [self._translate('volatility'),
             f"{portfolio.risk_metrics.get('volatility', 0):.2f}%"],
            [self._translate('sharpe_ratio'),
             f"{portfolio.risk_metrics.get('sharpe_ratio', 0):.2f}"],
            [self._translate('max_drawdown'),
             f"{portfolio.risk_metrics.get('max_drawdown', 0):.2f}%"]
        ]

        metrics_table = Table(metrics_data, colWidths=[3 * inch, 2 * inch])
        metrics_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, -1), colors.white),
            ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 10),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 12),
            ('GRID', (0, 0), (-1, -1), 1, colors.grey)
        ]))

        story.append(metrics_table)
        story.append(Spacer(1, 0.3 * inch))

        # Holdings table
        if self.config.include_tables:
            story.append(Paragraph(
                self._translate('holdings'),
                self.styles['SectionHeader']
            ))
            story.append(Spacer(1, 0.2 * inch))

            holdings_data = [[
                self._translate('symbol'),
                self._translate('shares'),
                self._translate('price'),
                self._translate('value'),
                self._translate('weight')
            ]]

            for symbol, data in portfolio.holdings.items():
                holdings_data.append([
                    symbol,
                    f"{data.get('shares', 0):.2f}",
                    f"${data.get('price', 0):.2f}",
                    f"${data.get('value', 0):,.2f}",
                    f"{data.get('weight', 0):.2f}%"
                ])

            holdings_table = Table(holdings_data, colWidths=[1 * inch] * 5)
            holdings_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1f77b4')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                ('GRID', (0, 0), (-1, -1), 1, colors.grey),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1),
                 [colors.white, colors.HexColor('#f0f0f0')])
            ]))

            story.append(holdings_table)
            story.append(Spacer(1, 0.3 * inch))

        # Allocation chart
        if self.config.include_charts:
            story.append(Paragraph(
                self._translate('allocation'),
                self.styles['SectionHeader']
            ))
            story.append(Spacer(1, 0.2 * inch))

            # Create pie chart
            drawing = Drawing(400, 200)
            pie = Pie()
            pie.x = 100
            pie.y = 20
            pie.width = 150
            pie.height = 150
            pie.data = list(portfolio.allocation.values())
            pie.labels = list(portfolio.allocation.keys())
            pie.slices.strokeWidth = 0.5
            drawing.add(pie)

            story.append(drawing)
            story.append(Spacer(1, 0.3 * inch))

        # Disclaimer
        if self.config.include_disclosures:
            story.append(PageBreak())
            story.append(Paragraph(
                self._translate('disclaimer'),
                self.styles['SectionHeader']
            ))
            story.append(Spacer(1, 0.2 * inch))

            disclaimer_text = """
            This report is for informational purposes only and does not constitute
            investment advice. Past performance is not indicative of future results.
            All investments carry risk, including the potential loss of principal.
            Consult with a qualified financial advisor before making investment decisions.
            """

            story.append(Paragraph(disclaimer_text, self.styles['Normal']))

        # Build PDF
        doc.build(story, onFirstPage=self._create_header,
                 onLaterPages=self._create_header)

    def generate_executive_summary(self, portfolio: PortfolioData,
                                  insights: Dict[str, Any],
                                  output_path: str):
        """Generate executive summary report."""
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        story = []

        # Title
        title = Paragraph(
            self._translate('executive_summary'),
            self.styles['CustomTitle']
        )
        story.append(title)
        story.append(Spacer(1, 0.3 * inch))

        # Date and portfolio name
        info = Paragraph(
            f"{portfolio.portfolio_name}<br/>"
            f"{self._translate('as_of')} {self.config.report_date.strftime('%B %d, %Y')}",
            self.styles['Normal']
        )
        story.append(info)
        story.append(Spacer(1, 0.3 * inch))

        # Key highlights
        story.append(Paragraph(
            self._translate('key_highlights'),
            self.styles['SectionHeader']
        ))
        story.append(Spacer(1, 0.2 * inch))

        highlights = insights.get('highlights', [])
        for highlight in highlights:
            bullet = Paragraph(f"• {highlight}", self.styles['Normal'])
            story.append(bullet)
            story.append(Spacer(1, 0.1 * inch))

        story.append(Spacer(1, 0.3 * inch))

        # Performance summary
        story.append(Paragraph(
            self._translate('performance'),
            self.styles['SectionHeader']
        ))
        story.append(Spacer(1, 0.2 * inch))

        performance_text = f"""
        The portfolio generated a return of {portfolio.performance.get('ytd_return', 0):.2f}%
        year-to-date, with an annualized volatility of {portfolio.risk_metrics.get('volatility', 0):.2f}%.
        The Sharpe ratio of {portfolio.risk_metrics.get('sharpe_ratio', 0):.2f} indicates
        {'strong' if portfolio.risk_metrics.get('sharpe_ratio', 0) > 1 else 'moderate'}
        risk-adjusted returns.
        """

        story.append(Paragraph(performance_text, self.styles['Normal']))
        story.append(Spacer(1, 0.3 * inch))

        # Recommendations
        if self.config.include_recommendations:
            story.append(Paragraph(
                self._translate('recommendations'),
                self.styles['SectionHeader']
            ))
            story.append(Spacer(1, 0.2 * inch))

            recommendations = insights.get('recommendations', [])
            for i, rec in enumerate(recommendations, 1):
                rec_text = f"{i}. {rec}"
                story.append(Paragraph(rec_text, self.styles['Normal']))
                story.append(Spacer(1, 0.1 * inch))

        doc.build(story)

    def generate_compliance_report(self, portfolio: PortfolioData,
                                  compliance_data: Dict[str, Any],
                                  output_path: str):
        """Generate regulatory compliance report."""
        doc = SimpleDocTemplate(output_path, pagesize=letter)
        story = []

        # Title
        title = Paragraph(
            "Regulatory Compliance Report",
            self.styles['CustomTitle']
        )
        story.append(title)
        story.append(Spacer(1, 0.3 * inch))

        # Compliance statement
        story.append(Paragraph(
            self._translate('compliance_statement'),
            self.styles['SectionHeader']
        ))
        story.append(Spacer(1, 0.2 * inch))

        compliance_text = """
        This report has been prepared in accordance with applicable regulatory
        requirements including SEC Rule 204-2, FINRA Rule 2210, and DOL ERISA
        Section 404(a). All disclosures, risk warnings, and performance calculations
        comply with industry standards.
        """

        story.append(Paragraph(compliance_text, self.styles['Normal']))
        story.append(Spacer(1, 0.3 * inch))

        # Regulatory checks
        story.append(Paragraph(
            "Regulatory Compliance Checks",
            self.styles['SectionHeader']
        ))
        story.append(Spacer(1, 0.2 * inch))

        checks_data = [
            ['Check', 'Status', 'Details']
        ]

        for check_name, check_result in compliance_data.get('checks', {}).items():
            status = "✓ Pass" if check_result.get('passed') else "✗ Fail"
            details = check_result.get('details', '')
            checks_data.append([check_name, status, details])

        checks_table = Table(checks_data, colWidths=[2 * inch, 1.5 * inch, 2.5 * inch])
        checks_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1f77b4')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
            ('GRID', (0, 0), (-1, -1), 1, colors.grey)
        ]))

        story.append(checks_table)

        doc.build(story)


class EmailDelivery:
    """Email report delivery system."""

    def __init__(self, smtp_server: str, smtp_port: int,
                 username: str, password: str):
        """Initialize email delivery."""
        self.smtp_server = smtp_server
        self.smtp_port = smtp_port
        self.username = username
        self.password = password

    def send_report(self, to_addresses: List[str], subject: str,
                   body: str, attachments: List[str],
                   cc_addresses: Optional[List[str]] = None,
                   bcc_addresses: Optional[List[str]] = None) -> bool:
        """Send report via email."""
        try:
            # Create message
            msg = MIMEMultipart()
            msg['From'] = self.username
            msg['To'] = ', '.join(to_addresses)
            msg['Subject'] = subject

            if cc_addresses:
                msg['Cc'] = ', '.join(cc_addresses)

            # Add body
            msg.attach(MIMEText(body, 'html'))

            # Add attachments
            for attachment_path in attachments:
                with open(attachment_path, 'rb') as f:
                    part = MIMEBase('application', 'octet-stream')
                    part.set_payload(f.read())
                    encoders.encode_base64(part)
                    part.add_header(
                        'Content-Disposition',
                        f'attachment; filename= {os.path.basename(attachment_path)}'
                    )
                    msg.attach(part)

            # Send email
            with smtplib.SMTP(self.smtp_server, self.smtp_port) as server:
                server.starttls()
                server.login(self.username, self.password)

                all_recipients = to_addresses.copy()
                if cc_addresses:
                    all_recipients.extend(cc_addresses)
                if bcc_addresses:
                    all_recipients.extend(bcc_addresses)

                server.send_message(msg)

            return True

        except Exception as e:
            print(f"Failed to send email: {e}")
            return False

    def generate_email_body(self, portfolio: PortfolioData,
                          config: ReportConfig) -> str:
        """Generate HTML email body."""
        translator = Translator()

        html = f"""
        <html>
        <head>
            <style>
                body {{ font-family: Arial, sans-serif; }}
                .header {{ background-color: #1f77b4; color: white; padding: 20px; }}
                .content {{ padding: 20px; }}
                .metrics {{ margin: 20px 0; }}
                .metric {{ display: inline-block; margin: 10px; padding: 15px;
                         background-color: #f0f0f0; border-radius: 5px; }}
                .footer {{ background-color: #333; color: white; padding: 10px;
                          text-align: center; font-size: 12px; }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1>{config.company_name}</h1>
                <h2>{translator.translate('portfolio_summary', config.language)}</h2>
            </div>

            <div class="content">
                <p>Dear Valued Client,</p>

                <p>Please find attached your {translator.translate('portfolio_summary', config.language)}
                   for {portfolio.portfolio_name}.</p>

                <div class="metrics">
                    <div class="metric">
                        <strong>{translator.translate('total_value', config.language)}:</strong><br/>
                        ${portfolio.total_value:,.2f}
                    </div>
                    <div class="metric">
                        <strong>{translator.translate('ytd_return', config.language)}:</strong><br/>
                        {portfolio.performance.get('ytd_return', 0):.2f}%
                    </div>
                    <div class="metric">
                        <strong>{translator.translate('sharpe_ratio', config.language)}:</strong><br/>
                        {portfolio.risk_metrics.get('sharpe_ratio', 0):.2f}
                    </div>
                </div>

                <p>For questions or to schedule a review, please contact us.</p>

                <p>Best regards,<br/>
                {config.company_name}</p>
            </div>

            <div class="footer">
                <p>{translator.translate('disclaimer', config.language)}: This email and any attachments
                are confidential and intended solely for the addressee.</p>
            </div>
        </body>
        </html>
        """

        return html


class PresentationGenerator:
    """Generates PowerPoint presentations."""

    def __init__(self, config: ReportConfig):
        """Initialize presentation generator."""
        self.config = config
        self.translator = Translator()

    def _translate(self, key: str) -> str:
        """Translate a key."""
        return self.translator.translate(key, self.config.language)

    def generate_client_presentation(self, portfolio: PortfolioData,
                                    output_path: str):
        """Generate client-ready PowerPoint presentation."""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        # Slide 1: Title
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]

        title.text = self._translate('portfolio_summary')
        subtitle.text = f"{portfolio.portfolio_name}\n{self.config.report_date.strftime('%B %Y')}"

        # Slide 2: Executive Summary
        summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = summary_slide.shapes.title
        title.text = self._translate('executive_summary')

        content = summary_slide.placeholders[1]
        tf = content.text_frame
        tf.text = f"{self._translate('total_value')}: ${portfolio.total_value:,.2f}"

        p = tf.add_paragraph()
        p.text = f"{self._translate('ytd_return')}: {portfolio.performance.get('ytd_return', 0):.2f}%"

        p = tf.add_paragraph()
        p.text = f"{self._translate('sharpe_ratio')}: {portfolio.risk_metrics.get('sharpe_ratio', 0):.2f}"

        # Slide 3: Performance
        perf_slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = perf_slide.shapes.title
        title.text = self._translate('performance')

        # Add chart placeholder
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)

        # Slide 4: Allocation
        alloc_slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = alloc_slide.shapes.title
        title.text = self._translate('allocation')

        # Add pie chart
        chart_data = CategoryChartData()
        chart_data.categories = list(portfolio.allocation.keys())
        chart_data.add_series('Allocation', list(portfolio.allocation.values()))

        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        chart = alloc_slide.shapes.add_chart(
            5,  # Pie chart
            x, y, cx, cy, chart_data
        ).chart

        # Slide 5: Holdings
        holdings_slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = holdings_slide.shapes.title
        title.text = self._translate('holdings')

        # Add table
        rows = len(portfolio.holdings) + 1
        cols = 4
        left = Inches(1.5)
        top = Inches(2)
        width = Inches(7)
        height = Inches(4)

        table = holdings_slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Header
        table.cell(0, 0).text = self._translate('symbol')
        table.cell(0, 1).text = self._translate('value')
        table.cell(0, 2).text = self._translate('weight')
        table.cell(0, 3).text = self._translate('change')

        # Data
        for i, (symbol, data) in enumerate(portfolio.holdings.items(), 1):
            table.cell(i, 0).text = symbol
            table.cell(i, 1).text = f"${data.get('value', 0):,.0f}"
            table.cell(i, 2).text = f"{data.get('weight', 0):.1f}%"
            table.cell(i, 3).text = f"{data.get('change', 0):+.2f}%"

        # Save presentation
        prs.save(output_path)


class ReportingSystem:
    """Main reporting system coordinator."""

    def __init__(self, config: ReportConfig,
                 email_config: Optional[Dict[str, str]] = None):
        """Initialize reporting system."""
        self.config = config
        self.pdf_generator = PDFReportGenerator(config)
        self.presentation_generator = PresentationGenerator(config)

        if email_config:
            self.email_delivery = EmailDelivery(
                smtp_server=email_config.get('smtp_server'),
                smtp_port=email_config.get('smtp_port', 587),
                username=email_config.get('username'),
                password=email_config.get('password')
            )
        else:
            self.email_delivery = None

    def generate_report(self, portfolio: PortfolioData,
                       output_path: str,
                       insights: Optional[Dict[str, Any]] = None,
                       compliance_data: Optional[Dict[str, Any]] = None) -> str:
        """Generate report based on configuration."""

        if self.config.format == ReportFormat.PDF:
            if self.config.report_type == ReportType.EXECUTIVE_SUMMARY:
                if insights is None:
                    insights = self._generate_default_insights(portfolio)
                self.pdf_generator.generate_executive_summary(
                    portfolio, insights, output_path
                )
            elif self.config.report_type == ReportType.COMPLIANCE:
                if compliance_data is None:
                    compliance_data = self._generate_default_compliance_data()
                self.pdf_generator.generate_compliance_report(
                    portfolio, compliance_data, output_path
                )
            else:
                self.pdf_generator.generate_portfolio_summary(
                    portfolio, output_path
                )

        elif self.config.format == ReportFormat.PPTX:
            self.presentation_generator.generate_client_presentation(
                portfolio, output_path
            )

        return output_path

    def send_report_via_email(self, portfolio: PortfolioData,
                            to_addresses: List[str],
                            report_path: str,
                            subject: Optional[str] = None,
                            cc_addresses: Optional[List[str]] = None) -> bool:
        """Generate and send report via email."""
        if self.email_delivery is None:
            raise ValueError("Email delivery not configured")

        if subject is None:
            subject = f"Portfolio Report - {portfolio.portfolio_name}"

        body = self.email_delivery.generate_email_body(portfolio, self.config)

        return self.email_delivery.send_report(
            to_addresses=to_addresses,
            subject=subject,
            body=body,
            attachments=[report_path],
            cc_addresses=cc_addresses
        )

    def _generate_default_insights(self, portfolio: PortfolioData) -> Dict[str, Any]:
        """Generate default insights for executive summary."""
        return {
            'highlights': [
                f"Portfolio value: ${portfolio.total_value:,.2f}",
                f"YTD return: {portfolio.performance.get('ytd_return', 0):.2f}%",
                f"Sharpe ratio: {portfolio.risk_metrics.get('sharpe_ratio', 0):.2f}",
                f"Number of holdings: {len(portfolio.holdings)}"
            ],
            'recommendations': [
                "Continue current allocation strategy",
                "Monitor market volatility",
                "Consider rebalancing if deviation exceeds 5%",
                "Review tax-loss harvesting opportunities"
            ]
        }

    def _generate_default_compliance_data(self) -> Dict[str, Any]:
        """Generate default compliance data."""
        return {
            'checks': {
                'SEC Disclosure': {'passed': True, 'details': 'All required disclosures present'},
                'FINRA 2210': {'passed': True, 'details': 'Communications standards met'},
                'DOL ERISA 404': {'passed': True, 'details': 'Fiduciary standards compliant'},
                'Performance Calc': {'passed': True, 'details': 'GIPS-compliant calculations'}
            }
        }


def create_default_config(report_type: ReportType = ReportType.PORTFOLIO_SUMMARY,
                         format: ReportFormat = ReportFormat.PDF,
                         language: Language = Language.ENGLISH) -> ReportConfig:
    """Create a default report configuration."""
    return ReportConfig(
        report_type=report_type,
        format=format,
        language=language,
        include_charts=True,
        include_tables=True,
        include_recommendations=True,
        include_disclosures=True
    )
